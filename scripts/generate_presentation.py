from pathlib import Path
from typing import Dict, Iterable, Mapping
from urllib.parse import quote_plus

import requests
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ASSETS_DIR = Path("assets/presentation")


def rgb_from_hex(hex_color: str) -> RGBColor:
    """Convert a CSS-style hex string to an RGBColor object."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def generate_slide_images(prompts: Mapping[str, str]) -> Dict[str, Path]:
    """Fetch AI-generated illustrations for each slide, with fallback placeholders."""
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    font = ImageFont.load_default()

    images: Dict[str, Path] = {}
    for title, prompt in prompts.items():
        safe_name = "".join(ch if ch.isalnum() else "_" for ch in title.lower())
        file_path = ASSETS_DIR / f"{safe_name}.png"

        if not file_path.exists():
            success = fetch_ai_image(prompt, file_path)
            if not success:
                create_fallback_image(title, font, file_path)

        images[title] = file_path

    return images


def fetch_ai_image(prompt: str, destination: Path) -> bool:
    """Attempt to download an AI-generated image using the Pollinations API."""
    try:
        query = quote_plus(prompt)
        url = f"https://image.pollinations.ai/prompt/{query}?width=800&height=450&nologo=1"
        response = requests.get(url, timeout=60)
        response.raise_for_status()
        destination.write_bytes(response.content)
        return True
    except Exception:
        return False


def create_fallback_image(title: str, font: ImageFont.FreeTypeFont, destination: Path) -> None:
    """Create a simple placeholder if the AI image request fails."""
    img = Image.new("RGB", (800, 450), color=(102, 90, 142))
    draw = ImageDraw.Draw(img)
    bbox = draw.textbbox((0, 0), title, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    draw.rectangle([(40, 40), (760, 410)], outline=(210, 199, 229), width=6)
    draw.text(
        ((800 - text_width) / 2, (450 - text_height) / 2),
        title,
        fill=(255, 255, 255),
        font=font,
    )
    img.save(destination)


def add_title_slide(prs: Presentation, title: str, subtitle: str, image_path: Path) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_background_color(slide, TEMPLATE_COLOR)

    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    add_image_to_slide(slide, image_path, prs)


def add_bullet_slide(prs: Presentation, heading: str, bullets: list[str], image_path: Path) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    set_background_color(slide, TEMPLATE_COLOR)

    title_shape = slide.shapes.title
    title_shape.text = heading
    title_shape.text_frame.paragraphs[0].font.bold = True

    body_shape = slide.placeholders[1]
    body_shape.left = Inches(0.7)
    body_shape.top = Inches(1.8)
    body_shape.width = prs.slide_width - Inches(4.5)
    body_shape.height = prs.slide_height - Inches(2.5)

    tf = body_shape.text_frame
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = bullet
        p.space_after = Pt(4)

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    add_image_to_slide(slide, image_path, prs)


def add_takeaway_slide(prs: Presentation, heading: str, bullets: list[str], image_path: Path) -> None:
    add_bullet_slide(prs, heading, bullets, image_path)


def add_image_to_slide(slide, image_path: Path, prs: Presentation) -> None:
    image_width = Inches(3.5)
    left = prs.slide_width - image_width - Inches(0.6)
    top = Inches(1.6)
    slide.shapes.add_picture(str(image_path), left, top, width=image_width)


def set_background_color(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def build_presentation() -> Presentation:
    prs = Presentation()

    slide_prompts = {
        "Corporate Governance and Legal Responsibilities of Engineers": "futuristic corporate boardroom with engineers collaborating, digital governance dashboard, purple accent lighting, cinematic concept art",
        "Governance Overview": "illustration of a balanced corporate governance wheel with diverse engineers and stakeholders, modern flat design, soft purple palette",
        "Governance Structures": "organizational chart visual with board, committees, and engineering team connections, isometric style, violet tones",
        "Policies Shaping Engineering": "engineer reviewing safety policy documents and sustainability checklist in a modern workspace, digital art",
        "Legal Duties of Engineers": "engineer signing compliance documents with legal scales and blueprint overlay, photo-realistic, professional lighting",
        "Ethics and Professional Standards": "diverse engineers pledging integrity with abstract ethical compass backdrop, semi-realistic digital painting",
        "Risk Management Role": "engineers analyzing risk dashboard with holographic warnings and mitigation icons, cyberpunk-inspired scene",
        "Partnering with Legal & Compliance": "engineer collaborating with legal advisor over contract documents, high-resolution illustration, warm tones",
        "Accountability & Documentation": "digital archive room with holographic records and engineer updating logs, sci-fi yet professional aesthetic",
        "Future Governance Trends": "AI-powered cityscape with connected infrastructure and governance icons, visionary concept art, purple glow",
        "Key Takeaways & Next Steps": "team of engineers celebrating success with checklist and roadmap, modern vector illustration, optimistic lighting",
    }

    images = generate_slide_images(slide_prompts)

    add_title_slide(
        prs,
        "Corporate Governance and Legal Responsibilities of Engineers",
        "Protecting Stakeholders Through Ethical Engineering",
        images["Corporate Governance and Legal Responsibilities of Engineers"],
    )

    slides_content = [
        (
            "Governance Overview",
            [
                "Corporate governance sets direction, control, and accountability for organizations.",
                "Balances interests of shareholders, management, customers, regulators, and the public.",
                "Engineers translate governance strategy into safe, compliant technical solutions.",
            ],
        ),
        (
            "Governance Structures",
            [
                "Boards establish strategy, approve risk appetite, and monitor performance impacting engineering programs.",
                "Audit, risk, and compliance committees oversee controls and reporting tied to technical work.",
                "Defined reporting lines connect engineering leads with governance bodies for oversight.",
            ],
        ),
        (
            "Policies Shaping Engineering",
            [
                "Quality and safety policies define technical requirements, testing, and approvals.",
                "Sustainability commitments steer design choices toward ESG-aligned outcomes.",
                "Transparency mandates accurate data reporting and stakeholder communication.",
            ],
        ),
        (
            "Legal Duties of Engineers",
            [
                "Duty of care requires diligence in design, testing, and deployment to protect public safety.",
                "Compliance with applicable regulations: building codes, environmental statutes, industry standards.",
                "Thorough documentation and traceability demonstrate compliance and mitigate liability.",
            ],
        ),
        (
            "Ethics and Professional Standards",
            [
                "Adhere to codes from IEEE, NSPE, or local engineering bodies to guide conduct.",
                "Identify conflicts of interest early; escalate issues through governance channels.",
                "Balance innovation with societal responsibilities and long-term stakeholder trust.",
            ],
        ),
        (
            "Risk Management Role",
            [
                "Embed risk assessments in design reviews, FMEA, and technical audits.",
                "Track key risk indicators and mitigation plans within governance frameworks.",
                "Capture lessons learned from engineering failures to strengthen future controls.",
            ],
        ),
        (
            "Partnering with Legal & Compliance",
            [
                "Support regulatory filings, certifications, and audit responses with accurate data.",
                "Provide technical diligence during mergers, acquisitions, and vendor evaluations.",
                "Craft precise specifications, warranties, and liability clauses within contracts.",
            ],
        ),
        (
            "Accountability & Documentation",
            [
                "Maintain design records, change logs, and approvals in accessible repositories.",
                "Use RACI matrices to clarify roles in technical decisions and sign-offs.",
                "Document rationale in decision registers to stand up under investigation or litigation.",
            ],
        ),
        (
            "Future Governance Trends",
            [
                "AI governance, data privacy, and cybersecurity regulations expand engineering duties.",
                "ESG reporting raises expectations for lifecycle sustainability and transparency.",
                "Global harmonization of engineering standards demands continuous monitoring.",
            ],
        ),
    ]

    for heading, bullets in slides_content:
        add_bullet_slide(prs, heading, bullets, images[heading])

    add_takeaway_slide(
        prs,
        "Key Takeaways & Next Steps",
        [
            "Engineering governance blends technical diligence, legal compliance, and ethical leadership.",
            "Immediate actions: review policies, assess compliance gaps, schedule training refreshers.",
            "Encourage dialogue: invite questions and plan deeper dives into governance priorities.",
        ],
        images["Key Takeaways & Next Steps"],
    )

    return prs


def main() -> None:
    prs = build_presentation()
    output_path = Path("Corporate_Governance_and_Legal_Responsibilities.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path.resolve()}")


if __name__ == "__main__":
    TEMPLATE_COLOR = rgb_from_hex("#d1c7e5")
    main()

